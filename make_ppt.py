from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── 색상 팔레트 ──────────────────────────────────────────────
C_BG      = RGBColor(0x1A, 0x1A, 0x2E)   # 어두운 네이비
C_ACCENT  = RGBColor(0x16, 0x21, 0x3E)   # 패널 배경
C_BLUE    = RGBColor(0x0F, 0x3C, 0x96)   # 포인트 블루
C_CYAN    = RGBColor(0x00, 0xB4, 0xD8)   # 시안
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY    = RGBColor(0xB0, 0xBE, 0xC5)
C_YELLOW  = RGBColor(0xFF, 0xD6, 0x00)
C_GREEN   = RGBColor(0x06, 0xD6, 0xA0)
C_RED     = RGBColor(0xFF, 0x47, 0x6F)

W = prs.slide_width
H = prs.slide_height
blank_layout = prs.slide_layouts[6]   # 완전 빈 레이아웃

# ── 헬퍼 함수 ────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=Pt(18), bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def bg(slide, color=C_BG):
    add_rect(slide, 0, 0, W, H, fill_color=color)


def title_bar(slide, title_text, subtitle_text=""):
    # 상단 파란 바
    add_rect(slide, 0, 0, W, Inches(1.1), fill_color=C_BLUE)
    add_text(slide, title_text,
             Inches(0.4), Inches(0.12), W - Inches(0.8), Inches(0.65),
             size=Pt(28), bold=True, color=C_WHITE)
    if subtitle_text:
        add_text(slide, subtitle_text,
                 Inches(0.4), Inches(0.72), W - Inches(0.8), Inches(0.35),
                 size=Pt(14), color=C_CYAN)
    # 하단 얇은 줄
    add_rect(slide, 0, H - Inches(0.35), W, Inches(0.35), fill_color=C_ACCENT)
    add_text(slide, "온라인 칠판 시스템  |  고급프로그래밍 팀 프로젝트",
             Inches(0.4), H - Inches(0.32), W - Inches(0.8), Inches(0.28),
             size=Pt(11), color=C_GRAY)


def bullet_box(slide, items, x, y, w, h, title="", bg_color=C_ACCENT):
    add_rect(slide, x, y, w, h, fill_color=bg_color,
             line_color=C_BLUE, line_width=Pt(1.5))
    oy = y + Inches(0.15)
    if title:
        add_text(slide, title, x + Inches(0.2), oy, w - Inches(0.4), Inches(0.38),
                 size=Pt(15), bold=True, color=C_CYAN)
        oy += Inches(0.38)
    for item in items:
        add_text(slide, item, x + Inches(0.25), oy, w - Inches(0.5), Inches(0.38),
                 size=Pt(13), color=C_WHITE)
        oy += Inches(0.36)


def code_box(slide, code_lines, x, y, w, h, title=""):
    add_rect(slide, x, y, w, h, fill_color=RGBColor(0x0D, 0x0D, 0x1A),
             line_color=C_CYAN, line_width=Pt(1))
    oy = y + Inches(0.12)
    if title:
        add_text(slide, title, x + Inches(0.2), oy, w - Inches(0.4), Inches(0.33),
                 size=Pt(12), bold=True, color=C_CYAN)
        oy += Inches(0.33)
    for line in code_lines:
        color = C_GRAY
        if line.strip().startswith("//") or line.strip().startswith("#"):
            color = RGBColor(0x6A, 0x99, 0x55)
        elif any(k in line for k in ["class ", "void ", "int ", "return", "public", "private", "struct"]):
            color = RGBColor(0x56, 0x9C, 0xD6)
        add_text(slide, line, x + Inches(0.2), oy, w - Inches(0.4), Inches(0.28),
                 size=Pt(10.5), color=color)
        oy += Inches(0.27)


# ════════════════════════════════════════════════════════════
#  슬라이드 1 — 표지
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s)
# 배경 장식 원
for cx, cy, r, alpha in [
    (Inches(11.5), Inches(1.0), Inches(2.5), RGBColor(0x0F,0x3C,0x96)),
    (Inches(1.0),  Inches(6.5), Inches(1.8), RGBColor(0x0F,0x3C,0x96)),
]:
    sh = s.shapes.add_shape(9, cx - r, cy - r, r*2, r*2)  # oval
    sh.fill.solid(); sh.fill.fore_color.rgb = alpha
    sh.line.fill.background()

add_rect(s, 0, Inches(2.8), Inches(0.15), Inches(1.9), fill_color=C_CYAN)
add_text(s, "Online Canvas", Inches(0.5), Inches(2.6), Inches(10), Inches(1.2),
         size=Pt(52), bold=True, color=C_WHITE)
add_text(s, "실시간 공유 온라인 칠판 시스템",
         Inches(0.5), Inches(3.8), Inches(10), Inches(0.7),
         size=Pt(26), color=C_CYAN)
add_text(s, "Qt C++  ·  Java Server  ·  React Web",
         Inches(0.5), Inches(4.55), Inches(10), Inches(0.5),
         size=Pt(18), color=C_GRAY)
add_rect(s, Inches(0.5), Inches(5.3), Inches(5.5), Inches(0.9),
         fill_color=C_BLUE)
add_text(s, "고급프로그래밍  팀 프로젝트   |   2024",
         Inches(0.65), Inches(5.38), Inches(5.2), Inches(0.5),
         size=Pt(16), color=C_WHITE, bold=True)

# ════════════════════════════════════════════════════════════
#  슬라이드 2 — 목차
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "목차", "Table of Contents")

items_l = [
    "01   프로젝트 개요",
    "02   개발 동기 및 목표",
    "03   전체 시스템 아키텍처",
    "04   통신 프로토콜 설계",
    "05   Java 서버 — MainServer",
    "06   Java 서버 — ClientHandler",
    "07   Java 서버 — WebServer (WebSocket)",
    "08   데이터 브로드캐스트 흐름",
]
items_r = [
    "09   C++ 클라이언트 구조",
    "10   Canvas 클래스 (그리기·지우기)",
    "11   MainWindow & 네트워크 처리",
    "12   직렬화 / 역직렬화",
    "13   PDF 수신 및 렌더링",
    "14   React 웹 프론트엔드",
    "15   개발 버전 히스토리",
    "16   결론 및 Q&A",
]
for i, item in enumerate(items_l):
    add_rect(s, Inches(0.4), Inches(1.3 + i*0.71), Inches(6.0), Inches(0.62),
             fill_color=C_ACCENT, line_color=C_BLUE, line_width=Pt(1))
    add_text(s, item, Inches(0.65), Inches(1.35 + i*0.71), Inches(5.7), Inches(0.5),
             size=Pt(14), color=C_WHITE)
for i, item in enumerate(items_r):
    add_rect(s, Inches(6.8), Inches(1.3 + i*0.71), Inches(6.0), Inches(0.62),
             fill_color=C_ACCENT, line_color=C_CYAN, line_width=Pt(1))
    add_text(s, item, Inches(7.05), Inches(1.35 + i*0.71), Inches(5.7), Inches(0.5),
             size=Pt(14), color=C_WHITE)

# ════════════════════════════════════════════════════════════
#  슬라이드 3 — 프로젝트 개요
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "01   프로젝트 개요", "Project Overview")

add_text(s, "Online Canvas",
         Inches(0.4), Inches(1.3), Inches(12.5), Inches(0.6),
         size=Pt(28), bold=True, color=C_CYAN)
add_text(s,
    "교수가 웹 브라우저에서 PDF를 업로드하면,\n"
    "학생들의 Qt C++ 클라이언트에 PDF가 실시간으로 전송되고\n"
    "모든 참여자가 같은 화면 위에 동시에 필기할 수 있는 시스템",
    Inches(0.4), Inches(2.0), Inches(12.5), Inches(1.3),
    size=Pt(17), color=C_WHITE)

features = [
    ("PDF 공유", "교수가 웹에서 PDF 업로드 → 학생 화면 자동 전송"),
    ("실시간 필기", "펜/지우개로 그린 획을 모든 클라이언트에 브로드캐스트"),
    ("페이지 제어", "교수가 페이지를 넘기면 학생 화면도 동기화"),
    ("다중 접속", "CopyOnWriteArrayList로 스레드 안전한 다중 클라이언트 관리"),
]
for i, (t, d) in enumerate(features):
    x = Inches(0.4 + (i % 2) * 6.45)
    y = Inches(3.5 + (i // 2) * 1.35)
    add_rect(s, x, y, Inches(6.1), Inches(1.15), fill_color=C_ACCENT,
             line_color=C_CYAN if i % 2 == 0 else C_BLUE, line_width=Pt(1.5))
    add_text(s, t, x + Inches(0.2), y + Inches(0.1), Inches(5.7), Inches(0.38),
             size=Pt(15), bold=True, color=C_CYAN)
    add_text(s, d, x + Inches(0.2), y + Inches(0.5), Inches(5.7), Inches(0.55),
             size=Pt(12.5), color=C_GRAY)

# ════════════════════════════════════════════════════════════
#  슬라이드 4 — 개발 동기
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "02   개발 동기 및 목표", "Motivation & Goals")

motivations = [
    "💡  기존 화상강의 도구는 교수 화면만 공유 — 학생이 필기 불가",
    "💡  별도 필기 앱과 강의 화면을 동시에 쓰면 화면 전환이 번거로움",
    "💡  같은 PDF 위에 교수·학생 모두가 함께 필기하는 협업 도구 필요",
]
goals = [
    "✅  Qt C++ 그래픽 프로그래밍 실습",
    "✅  TCP 소켓 통신 & 바이너리 직렬화 학습",
    "✅  Java WebSocket 서버 설계 경험",
    "✅  React 웹과 네이티브 앱 간 실시간 연동",
]

add_text(s, "개발 동기", Inches(0.4), Inches(1.25), Inches(5.8), Inches(0.45),
         size=Pt(17), bold=True, color=C_CYAN)
for i, m in enumerate(motivations):
    add_rect(s, Inches(0.4), Inches(1.75 + i * 0.95), Inches(5.9), Inches(0.8),
             fill_color=C_ACCENT, line_color=C_BLUE, line_width=Pt(1))
    add_text(s, m, Inches(0.6), Inches(1.8 + i * 0.95), Inches(5.5), Inches(0.65),
             size=Pt(13), color=C_WHITE)

add_text(s, "개발 목표", Inches(7.0), Inches(1.25), Inches(5.8), Inches(0.45),
         size=Pt(17), bold=True, color=C_GREEN)
for i, g in enumerate(goals):
    add_rect(s, Inches(7.0), Inches(1.75 + i * 0.95), Inches(5.9), Inches(0.8),
             fill_color=C_ACCENT, line_color=C_GREEN, line_width=Pt(1))
    add_text(s, g, Inches(7.2), Inches(1.8 + i * 0.95), Inches(5.5), Inches(0.65),
             size=Pt(13), color=C_WHITE)

add_rect(s, Inches(6.5), Inches(1.2), Inches(0.05), Inches(5.5), fill_color=C_GRAY)

# ════════════════════════════════════════════════════════════
#  슬라이드 5 — 전체 아키텍처
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "03   전체 시스템 아키텍처", "System Architecture")

# 박스 3개
boxes = [
    (Inches(0.35), Inches(1.25), Inches(3.7), Inches(5.5), C_BLUE,   "교수 웹 브라우저\n(React + Vite)"),
    (Inches(4.8),  Inches(1.25), Inches(3.7), Inches(5.5), C_ACCENT, "Java 서버"),
    (Inches(9.25), Inches(1.25), Inches(3.7), Inches(5.5), C_BLUE,   "Qt C++ 클라이언트\n(학생 × N)"),
]
for bx, by, bw, bh, bc, bt in boxes:
    add_rect(s, bx, by, bw, bh, fill_color=bc,
             line_color=C_CYAN, line_width=Pt(2))
    add_text(s, bt, bx + Inches(0.1), by + Inches(0.15), bw - Inches(0.2), Inches(0.75),
             size=Pt(15), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 웹 내부
web_items = ["StartScreen", "UploadScreen", "LectureMode", "PDF 업로드", "페이지 전환 버튼"]
for i, wi in enumerate(web_items):
    add_rect(s, Inches(0.55), Inches(2.2 + i*0.68), Inches(3.3), Inches(0.55),
             fill_color=RGBColor(0x0F,0x3C,0x96), line_color=C_WHITE, line_width=Pt(0.5))
    add_text(s, wi, Inches(0.65), Inches(2.25 + i*0.68), Inches(3.1), Inches(0.42),
             size=Pt(12), color=C_WHITE, align=PP_ALIGN.CENTER)

# 서버 내부
srv_items = ["MainServer (TCP :8081)", "WebServer (WS :8082)", "ClientHandler × N", "브로드캐스트 로직"]
for i, si in enumerate(srv_items):
    add_rect(s, Inches(5.0), Inches(2.2 + i*0.8), Inches(3.3), Inches(0.62),
             fill_color=RGBColor(0x0D,0x22,0x3A), line_color=C_CYAN, line_width=Pt(0.5))
    add_text(s, si, Inches(5.1), Inches(2.25 + i*0.8), Inches(3.1), Inches(0.5),
             size=Pt(12), color=C_CYAN, align=PP_ALIGN.CENTER)

# Qt 내부
qt_items = ["MainWindow (UI + 소켓)", "Canvas (그리기 위젯)", "QPdfDocument", "QTcpSocket"]
for i, qi in enumerate(qt_items):
    add_rect(s, Inches(9.45), Inches(2.2 + i*0.8), Inches(3.3), Inches(0.62),
             fill_color=RGBColor(0x0F,0x3C,0x96), line_color=C_WHITE, line_width=Pt(0.5))
    add_text(s, qi, Inches(9.55), Inches(2.25 + i*0.8), Inches(3.1), Inches(0.5),
             size=Pt(12), color=C_WHITE, align=PP_ALIGN.CENTER)

# 화살표 텍스트
add_text(s, "WebSocket\n(8082)", Inches(4.1), Inches(3.3), Inches(0.6), Inches(0.7),
         size=Pt(9), color=C_YELLOW, align=PP_ALIGN.CENTER)
add_text(s, "TCP\n(8081)", Inches(8.5), Inches(3.3), Inches(0.6), Inches(0.7),
         size=Pt(9), color=C_YELLOW, align=PP_ALIGN.CENTER)
add_text(s, "◀──▶", Inches(4.05), Inches(3.9), Inches(0.7), Inches(0.4),
         size=Pt(14), color=C_YELLOW, align=PP_ALIGN.CENTER)
add_text(s, "◀──▶", Inches(8.5), Inches(3.9), Inches(0.7), Inches(0.4),
         size=Pt(14), color=C_YELLOW, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
#  슬라이드 6 — 통신 프로토콜
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "04   통신 프로토콜 설계", "Binary Protocol — BigEndian")

add_text(s, "모든 통신은 BigEndian 바이너리 포맷으로 통일 (TCP & WebSocket)",
         Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.4),
         size=Pt(14), color=C_GRAY)

rows = [
    ("type", "설명", "패킷 구조"),
    ("0  (PEN)",   "펜 획 전송",    "int32(0) | int64 strokeId | int32 N | int32 color | int32 size | float32 x,y × N"),
    ("1  (ERASE)", "획 지우기",      "int32(1) | int64 strokeId"),
    ("2  (PDF)",   "PDF 파일 전송", "int32(2) | int32 length | byte[] × length"),
    ("3  (PAGE)",  "페이지 전환",   "int32(3) | int32 pageNum  (1-based)"),
]
colors = [C_BLUE, C_ACCENT, C_ACCENT, C_ACCENT, C_ACCENT]
tcols = [C_WHITE, C_CYAN, C_GREEN, C_YELLOW, RGBColor(0xFF,0x99,0x44)]

for i, (t, d, p) in enumerate(rows):
    y = Inches(1.75 + i * 0.98)
    add_rect(s, Inches(0.4), y, Inches(1.5), Inches(0.82),
             fill_color=colors[i], line_color=C_CYAN, line_width=Pt(1))
    add_rect(s, Inches(1.9), y, Inches(2.5), Inches(0.82),
             fill_color=colors[i], line_color=C_CYAN, line_width=Pt(1))
    add_rect(s, Inches(4.4), y, Inches(8.5), Inches(0.82),
             fill_color=colors[i], line_color=C_CYAN, line_width=Pt(1))
    add_text(s, t, Inches(0.5), y + Inches(0.15), Inches(1.3), Inches(0.55),
             size=Pt(13), bold=(i==0), color=tcols[i])
    add_text(s, d, Inches(2.0), y + Inches(0.15), Inches(2.2), Inches(0.55),
             size=Pt(13), bold=(i==0), color=tcols[i])
    add_text(s, p, Inches(4.55), y + Inches(0.15), Inches(8.1), Inches(0.55),
             size=Pt(11.5 if i>0 else 13), bold=(i==0), color=tcols[i])

add_text(s, "※ 좌표는 [0, 1] 정규화된 float32 — 해상도 무관 재생 가능",
         Inches(0.4), Inches(6.7), Inches(12.5), Inches(0.4),
         size=Pt(12), italic=True, color=C_GRAY)

# ════════════════════════════════════════════════════════════
#  슬라이드 7 — MainServer.java
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "05   Java 서버 — MainServer", "Entry Point & Client Registry")

bullet_box(s,
    ["TCP 포트 8081에서 C++ 클라이언트 수신 대기",
     "WebServer(8082) 동시 기동",
     "접속 시 10바이트 사용자 ID 수신",
     "HTTP 브라우저 오접속 차단 (GET/HTTP 감지)",
     "CopyOnWriteArrayList로 스레드 안전 클라이언트 명단 관리",
     "ClientHandler 스레드 생성 후 clients 리스트에 등록"],
    Inches(0.4), Inches(1.25), Inches(5.0), Inches(5.5),
    title="핵심 역할")

code_box(s,
    ["public class MainServer {",
     "  public static WebServer webserver",
     "            = new WebServer(8082);",
     "  public static CopyOnWriteArrayList",
     "    <ClientHandler> clients = new ...",
     "",
     "  public static void main(...) {",
     "    ServerSocket ss = new ServerSocket(8081);",
     "    webserver.start();",
     "    while (true) {",
     "      Socket sock = ss.accept();",
     "      // ID 읽기, HTTP 차단",
     "      ClientHandler h = new ClientHandler(..);",
     "      clients.add(h);",
     "      new Thread(h).start();",
     "    }",
     "  }",
     "}",
    ],
    Inches(5.7), Inches(1.25), Inches(7.2), Inches(5.5),
    title="MainServer.java")

# ════════════════════════════════════════════════════════════
#  슬라이드 8 — ClientHandler.java
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "06   Java 서버 — ClientHandler", "Per-Client Thread & Broadcast")

bullet_box(s,
    ["각 C++ 클라이언트마다 독립 스레드로 실행",
     "DataInputStream으로 type 필드 먼저 읽음",
     "type=0 (PEN): 점 데이터 읽어 재조립 후 전체 브로드캐스트",
     "type=1 (ERASE): strokeId 읽어 브로드캐스트",
     "자기 자신은 브로드캐스트에서 제외 (이중 그리기 방지)",
     "클라이언트 종료 시 clients 리스트에서 자동 제거",
     "WebServer에도 동시 broadcast 가능"],
    Inches(0.4), Inches(1.25), Inches(5.0), Inches(5.5),
    title="핵심 역할")

code_box(s,
    ["// type=0 PEN 수신 & 브로드캐스트",
     "long strokeId = dis.readLong();",
     "int  pointCnt = dis.readInt();",
     "int  color    = dis.readInt();",
     "int  penSize  = dis.readInt();",
     "",
     "ByteBuffer buf = ByteBuffer",
     "  .allocate(24 + pointCnt * 8);",
     "buf.putInt(type);",
     "buf.putLong(strokeId);",
     "// ... 점 데이터 채우기 ...",
     "",
     "for (ClientHandler c : clients) {",
     "  if (c != this)       // 자신 제외",
     "    c.sendData(buf.array());",
     "}",
    ],
    Inches(5.7), Inches(1.25), Inches(7.2), Inches(5.5),
    title="ClientHandler.java — 핵심 로직")

# ════════════════════════════════════════════════════════════
#  슬라이드 9 — WebServer.java
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "07   Java 서버 — WebServer", "WebSocket Bridge (Port 8082)")

bullet_box(s,
    ["java-websocket 라이브러리 기반",
     "교수 브라우저와 WebSocket으로 통신",
     "TCP와 동일한 바이너리 프로토콜 사용",
     "onMessage(ByteBuffer)로 바이너리 프레임 수신",
     "type 파싱 후 모든 C++ 클라이언트에 relay",
     "type 2 (PDF): 교수가 업로드한 PDF 전달",
     "type 3 (PAGE): 교수가 누른 페이지 번호 전달"],
    Inches(0.4), Inches(1.25), Inches(5.0), Inches(5.5),
    title="핵심 역할")

code_box(s,
    ["@Override",
     "public void onMessage(WebSocket conn,",
     "                      ByteBuffer msg) {",
     "  msg.order(ByteOrder.BIG_ENDIAN);",
     "  byte[] packet = ...; // 전체 복사",
     "  int type = msg.getInt();",
     "",
     "  switch (type) {",
     "    case MSG_PDF:  // type=2",
     "      int len = msg.getInt();",
     "      // PDF 바이트 로깅",
     "      break;",
     "    case MSG_PAGE: // type=3",
     "      int page = msg.getInt();",
     "      break;",
     "  }",
     "  // C++ 클라이언트 전원에게 relay",
     "  for (ClientHandler c : clients)",
     "    c.sendData(packet);",
     "}",
    ],
    Inches(5.7), Inches(1.25), Inches(7.2), Inches(5.5),
    title="WebServer.java — onMessage")

# ════════════════════════════════════════════════════════════
#  슬라이드 10 — 데이터 흐름
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "08   데이터 브로드캐스트 흐름", "Data Flow Diagram")

steps = [
    ("①", "교수 웹에서\nPDF 업로드",    Inches(0.3),  C_BLUE),
    ("②", "WebSocket으로\nJava 서버 전송", Inches(2.9),  C_CYAN),
    ("③", "Java 서버가\ntype=2 패킷 생성", Inches(5.5),  C_BLUE),
    ("④", "TCP로 C++\n클라이언트 전달",  Inches(8.1),  C_CYAN),
    ("⑤", "Canvas에 PDF\n배경 렌더링",   Inches(10.7), C_BLUE),
]
for num, desc, x, c in steps:
    add_rect(s, x, Inches(1.5), Inches(2.35), Inches(1.6),
             fill_color=c, line_color=C_WHITE, line_width=Pt(1))
    add_text(s, num, x + Inches(0.05), Inches(1.55), Inches(0.5), Inches(0.5),
             size=Pt(18), bold=True, color=C_YELLOW)
    add_text(s, desc, x + Inches(0.1), Inches(1.95), Inches(2.1), Inches(1.0),
             size=Pt(12), color=C_WHITE)

# 화살표
for i in range(4):
    add_text(s, "▶", Inches(2.55 + i*2.6), Inches(2.05), Inches(0.4), Inches(0.5),
             size=Pt(20), color=C_YELLOW, align=PP_ALIGN.CENTER)

# 펜 흐름
add_text(s, "펜 획 브로드캐스트 흐름", Inches(0.4), Inches(3.5), Inches(12.5), Inches(0.45),
         size=Pt(16), bold=True, color=C_GREEN)
pen_steps = [
    ("학생 A가\n마우스로 그림", Inches(0.3)),
    ("strokeFinished\n시그널 emit", Inches(2.9)),
    ("TCP로 서버에\n직렬화 전송",  Inches(5.5)),
    ("자신 제외\n전원 브로드캐스트", Inches(8.1)),
    ("학생 B,C\n화면에 획 표시",  Inches(10.7)),
]
for desc, x in pen_steps:
    add_rect(s, x, Inches(4.1), Inches(2.35), Inches(1.5),
             fill_color=RGBColor(0x06,0x4E,0x2A), line_color=C_GREEN, line_width=Pt(1))
    add_text(s, desc, x + Inches(0.1), Inches(4.2), Inches(2.1), Inches(1.2),
             size=Pt(12), color=C_WHITE)
for i in range(4):
    add_text(s, "▶", Inches(2.55 + i*2.6), Inches(4.65), Inches(0.4), Inches(0.5),
             size=Pt(20), color=C_GREEN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════
#  슬라이드 11 — C++ 클라이언트 구조
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "09   C++ 클라이언트 구조", "Qt Application Architecture")

add_text(s, "Qt 클래스 계층 구조", Inches(0.4), Inches(1.25), Inches(5.5), Inches(0.45),
         size=Pt(17), bold=True, color=C_CYAN)

tree_lines = [
    "QApplication",
    "  └── MainWindow  (QMainWindow)",
    "        ├── Ui::MainWindow  (생성된 UI)",
    "        │     ├── canvasContainer  (QWidget)",
    "        │     │     └── Canvas  (QWidget)",
    "        │     ├── penSlider  (QSlider)",
    "        │     ├── penSizeEdit  (QLineEdit)",
    "        │     └── colorPickerButton  (QPushButton)",
    "        └── QTcpSocket  (서버 통신)",
]
for i, line in enumerate(tree_lines):
    color = C_CYAN if "MainWindow" in line or "Canvas" in line else \
            C_GREEN if "QTcpSocket" in line else C_WHITE
    add_text(s, line, Inches(0.5), Inches(1.8 + i * 0.47), Inches(5.5), Inches(0.42),
             size=Pt(12.5), color=color)

bullet_box(s,
    ["main.cpp: QApplication + MainWindow 생성",
     "mainwindow.h/cpp: UI, 소켓, 직렬화 담당",
     "canvas.h/cpp: 실제 그리기 위젯",
     "mainwindow.ui: Qt Designer XML 레이아웃",
     "CMakeLists.txt: Qt6 Core/Gui/Widgets/Network/Pdf"],
    Inches(6.3), Inches(1.25), Inches(6.6), Inches(3.0),
    title="파일 구성")

bullet_box(s,
    ["Qt Signal/Slot로 이벤트 연결",
     "setFixedSize()로 캔버스 크기 고정",
     "QRandomGenerator로 획 ID 생성",
     "QDataStream으로 BigEndian 직렬화"],
    Inches(6.3), Inches(4.45), Inches(6.6), Inches(2.35),
    title="주요 기술 포인트")

# ════════════════════════════════════════════════════════════
#  슬라이드 12 — Canvas 클래스
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "10   Canvas 클래스", "Drawing Widget — canvas.h / canvas.cpp")

bullet_box(s,
    ["QWidget 상속, 마우스 이벤트 override",
     "strokes (QList<Stroke>): 완성된 획 저장",
     "currentStroke: 현재 그리는 중인 획",
     "Tool enum: Pen / Eraser",
     "penColor, penSize, eraserSize 속성",
     "QPdfDocument: PDF 문서 보유",
     "backgroundImage (QImage): PDF 래스터 캐시"],
    Inches(0.4), Inches(1.25), Inches(5.5), Inches(5.5),
    title="주요 멤버")

code_box(s,
    ["// 마우스 뗄 때 → 서버로 전송",
     "void Canvas::mouseReleaseEvent(...) {",
     "  if (tool == Pen && !stroke.empty()) {",
     "    strokes.append(currentStroke);",
     "    emit strokeFinished(currentStroke);",
     "  }",
     "}",
     "",
     "// 지우개: 획 rect 충돌 감지",
     "strokes.removeIf([&](const Stroke& s){",
     "  for (const QPoint& p : s.points)",
     "    if (eraserRect.contains(p)) {",
     "      emit eraseRequested(s.id);",
     "      return true;",
     "    }",
     "  return false;",
     "});",
    ],
    Inches(6.0), Inches(1.25), Inches(6.9), Inches(5.5),
    title="핵심 이벤트 로직")

# ════════════════════════════════════════════════════════════
#  슬라이드 13 — MainWindow 네트워크
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "11   MainWindow & 네트워크 처리", "TCP Socket & Deserialization")

bullet_box(s,
    ["QTcpSocket::readyRead → onSocketReadyRead()",
     "receiveBuffer에 데이터 누적 (스트리밍 대응)",
     "QDataStream으로 BigEndian 역직렬화",
     "type 필드로 분기: 0=PEN, 1=ERASE, 2=PDF, 3=PAGE",
     "불완전 패킷이면 break → 다음 readyRead 대기",
     "처리 완료 후 receiveBuffer.remove()로 소비",
    ],
    Inches(0.4), Inches(1.25), Inches(5.5), Inches(4.5),
    title="수신 처리 흐름")

code_box(s,
    ["void MainWindow::onSocketReadyRead() {",
     "  receiveBuffer += socket->readAll();",
     "  QDataStream ds(&receiveBuffer, ReadOnly);",
     "  ds.setByteOrder(BigEndian);",
     "",
     "  while (true) {",
     "    if (bytesAvailable < 4) break;",
     "    qint32 type; ds >> type;",
     "",
     "    if (type == 0) { // PEN",
     "      if (bytesAvail < 20) break;",
     "      // id, N, color, size, points 읽기",
     "      canvas->onStrokeReceived(stroke);",
     "    } else if (type == 2) { // PDF",
     "      // length 확인 → 전체 도착 대기",
     "      canvas->onPdfReceived(pdfData);",
     "    }",
     "    receiveBuffer.remove(0, consumed);",
     "  }",
     "}",
    ],
    Inches(6.0), Inches(1.25), Inches(6.9), Inches(5.5),
    title="onSocketReadyRead() 핵심 로직")

# ════════════════════════════════════════════════════════════
#  슬라이드 14 — 직렬화 / PDF 렌더링
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "12 · 13   직렬화 & PDF 렌더링", "Serialization & PDF Rendering")

# 좌: 직렬화
add_text(s, "좌표 정규화 직렬화", Inches(0.4), Inches(1.25), Inches(5.8), Inches(0.42),
         size=Pt(16), bold=True, color=C_CYAN)
code_box(s,
    ["// 송신: 픽셀 → [0,1] 정규화",
     "float nx = (float)p.x() / canvas->width();",
     "float ny = (float)p.y() / canvas->height();",
     "ds << nx << ny;",
     "",
     "// 수신: [0,1] → 픽셀 역정규화",
     "float x, y; ds >> x >> y;",
     "QPoint pt(",
     "  (int)(x * canvas->width()),",
     "  (int)(y * canvas->height())",
     ");",
    ],
    Inches(0.4), Inches(1.75), Inches(5.8), Inches(3.5),
    title="sendStroke() — 정규화 이유: 해상도 독립")

add_text(s, "✓ 서로 다른 화면 크기 클라이언트도 동일하게 표시",
         Inches(0.4), Inches(5.4), Inches(5.8), Inches(0.5),
         size=Pt(12), color=C_GREEN)

# 우: PDF 렌더링
add_text(s, "PDF 수신 및 렌더링", Inches(6.7), Inches(1.25), Inches(6.2), Inches(0.42),
         size=Pt(16), bold=True, color=C_YELLOW)
code_box(s,
    ["// QBuffer: 메모리에서 바로 로드",
     "QBuffer *buf = new QBuffer(this);",
     "buf->setData(pdfData);",
     "buf->open(QIODevice::ReadOnly);",
     "pdfDocument->load(buf);",
     "",
     "// 현재 페이지 래스터화",
     "void Canvas::renderCurrentPage() {",
     "  backgroundImage = pdfDocument->render(",
     "    currentPageIndex,",
     "    QSize(width(), height())",
     "  );",
     "}",
     "",
     "// paintEvent에서 배경으로 먼저 그림",
     "painter.drawImage(0, 0, backgroundImage);",
    ],
    Inches(6.7), Inches(1.75), Inches(6.2), Inches(5.0),
    title="onPdfReceived() & renderCurrentPage()")

# ════════════════════════════════════════════════════════════
#  슬라이드 15 — 웹 프론트엔드
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "14   React 웹 프론트엔드", "Professor's Web Interface")

# 화면 흐름
screens = ["StartScreen\n세션 코드 입력", "UploadScreen\nPDF 업로드", "LectureMode\n페이지 전환 · 필기"]
for i, sc in enumerate(screens):
    x = Inches(0.5 + i * 4.2)
    add_rect(s, x, Inches(1.3), Inches(3.7), Inches(1.6),
             fill_color=C_BLUE, line_color=C_CYAN, line_width=Pt(1.5))
    add_text(s, sc, x + Inches(0.1), Inches(1.4), Inches(3.5), Inches(1.3),
             size=Pt(14), color=C_WHITE, align=PP_ALIGN.CENTER)
    if i < 2:
        add_text(s, "▶", Inches(4.0 + i * 4.2), Inches(1.9), Inches(0.35), Inches(0.5),
                 size=Pt(18), color=C_YELLOW)

bullet_box(s,
    ["React + Vite 기반 SPA",
     "WebSocket (ws://서버IP:8082) 연결 유지",
     "ws.binaryType = 'arraybuffer' — 바이너리 수신",
     "PDF 업로드 → ArrayBuffer로 읽어 type=2 패킷 전송",
     "페이지 버튼 → type=3 패킷 전송",
     "WsStatus 컴포넌트로 연결 상태 실시간 표시"],
    Inches(0.4), Inches(3.2), Inches(5.9), Inches(3.9),
    title="기술 스택 & 주요 기능")

code_box(s,
    ["// WebSocket 바이너리 전송 예시",
     "const buf = new ArrayBuffer(8);",
     "const view = new DataView(buf);",
     "view.setInt32(0, 3);      // type=PAGE",
     "view.setInt32(4, pageNum);",
     "ws.send(buf);",
     "",
     "// PDF 전송",
     "const reader = new FileReader();",
     "reader.onload = (e) => {",
     "  const ab = e.target.result;",
     "  ws.send(buildPdfPacket(ab));",
     "};",
    ],
    Inches(6.6), Inches(3.2), Inches(6.4), Inches(3.9),
    title="웹 → 서버 전송 로직")

# ════════════════════════════════════════════════════════════
#  슬라이드 16 — 버전 히스토리 & 마무리
# ════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank_layout)
bg(s); title_bar(s, "15 · 16   개발 히스토리 & 마무리", "Version History & Conclusion")

# 버전 타임라인
versions = [
    ("v1.0", "그림판 기본",    ["프로그램 실행 시 윈도우 생성", "검정색 펜으로 그리기"],           C_BLUE),
    ("v1.1", "UI 확장",        ["지우개 추가", "메인 윈도우에 그림판 통합", "사이드 패널 구성"],    C_CYAN),
    ("v1.2", "네트워크 연동\n(진행중)", ["펜 색 변경", "그림 데이터 서버 전송", "PDF 수신 렌더링"], C_GREEN),
]
for i, (ver, title, items, c) in enumerate(versions):
    x = Inches(0.35 + i * 4.3)
    add_rect(s, x, Inches(1.3), Inches(4.0), Inches(3.6),
             fill_color=C_ACCENT, line_color=c, line_width=Pt(2))
    add_text(s, ver, x + Inches(0.15), Inches(1.4), Inches(1.2), Inches(0.45),
             size=Pt(16), bold=True, color=c)
    add_text(s, title, x + Inches(0.15), Inches(1.85), Inches(3.7), Inches(0.5),
             size=Pt(13), color=C_WHITE)
    for j, it in enumerate(items):
        add_text(s, "• " + it, x + Inches(0.2), Inches(2.45 + j * 0.48), Inches(3.6), Inches(0.42),
                 size=Pt(12), color=C_GRAY)

# 추가 예정
add_text(s, "추가 예정", Inches(0.35), Inches(5.1), Inches(12.6), Inches(0.42),
         size=Pt(15), bold=True, color=C_YELLOW)
todos = ["통신 최적화 (획 순서 보장)", "윈도우 전체 화면 & UI 개선", "세션 코드 기반 입장 시스템"]
for i, td in enumerate(todos):
    add_rect(s, Inches(0.35 + i * 4.3), Inches(5.6), Inches(4.0), Inches(0.7),
             fill_color=RGBColor(0x3A,0x2E,0x00), line_color=C_YELLOW, line_width=Pt(1))
    add_text(s, "• " + td, Inches(0.55 + i * 4.3), Inches(5.68), Inches(3.7), Inches(0.55),
             size=Pt(12.5), color=C_YELLOW)

add_rect(s, 0, Inches(6.4), W, Inches(0.75), fill_color=C_BLUE)
add_text(s, "감사합니다  |  Q & A",
         Inches(0.4), Inches(6.45), W - Inches(0.8), Inches(0.6),
         size=Pt(22), bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# ── 저장 ─────────────────────────────────────────────────────
out = "/Users/moonsunghyo/2학년/고프/C-_project/발표자료_OnlineCanvas.pptx"
prs.save(out)
print("저장 완료:", out)
